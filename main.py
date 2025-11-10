import os
from io import BytesIO
from datetime import datetime
from fastapi import FastAPI
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse

# PPTX imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
def read_root():
    return {"message": "Hello from FastAPI Backend!"}

@app.get("/api/hello")
def hello():
    return {"message": "Hello from the backend API!"}

@app.get("/test")
def test_database():
    """Test endpoint to check if database is available and accessible"""
    response = {
        "backend": "✅ Running",
        "database": "❌ Not Available",
        "database_url": None,
        "database_name": None,
        "connection_status": "Not Connected",
        "collections": []
    }
    try:
        from database import db
        if db is not None:
            response["database"] = "✅ Available"
            response["database_url"] = "✅ Configured"
            response["database_name"] = db.name if hasattr(db, 'name') else "✅ Connected"
            response["connection_status"] = "Connected"
            try:
                collections = db.list_collection_names()
                response["collections"] = collections[:10]
                response["database"] = "✅ Connected & Working"
            except Exception as e:
                response["database"] = f"⚠️  Connected but Error: {str(e)[:50]}"
        else:
            response["database"] = "⚠️  Available but not initialized"
    except ImportError:
        response["database"] = "❌ Database module not found (run enable-database first)"
    except Exception as e:
        response["database"] = f"❌ Error: {str(e)[:50]}"

    response["database_url"] = "✅ Set" if os.getenv("DATABASE_URL") else "❌ Not Set"
    response["database_name"] = "✅ Set" if os.getenv("DATABASE_NAME") else "❌ Not Set"
    return response


def _add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullets_slide(prs: Presentation, title: str, bullets: list[str]):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.level = 0


def _add_two_column_slide(prs: Presentation, title: str, left_title: str, left_items: list[str], right_title: str, right_items: list[str]):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(4.3), Inches(5))
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5))

    for title_text, items, box in [
        (left_title, left_items, left_box),
        (right_title, right_items, right_box)
    ]:
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.bold = True
        run.font.size = Pt(20)
        p.alignment = PP_ALIGN.LEFT
        # items
        for it in items:
            pi = tf.add_paragraph()
            pi.text = f"• {it}"
            pi.level = 0
            pi.font.size = Pt(16)


def build_ipb_ui_presentation() -> BytesIO:
    prs = Presentation()

    # Title
    _add_title_slide(
        prs,
        "Profil IPB University & Universitas Indonesia",
        "Ringkasan jalur masuk dan fakultas (disusun otomatis)"
    )

    # IPB Overview
    _add_bullets_slide(
        prs,
        "IPB University – Ringkasan",
        [
            "Perguruan tinggi negeri unggul berlokasi di Bogor, Jawa Barat",
            "Bertransformasi menjadi Techno-Socio Entrepreneurial University",
            "Keunggulan: biosains tropika, pertanian, kelautan, dan teknologi terkait",
            "Sering masuk Top 50 dunia untuk Pertanian & Kehutanan (QS WUR)",
            "Akreditasi: Unggul; Program: Vokasi hingga Pascasarjana",
            "Kampus utama di Dramaga; inovasi untuk kemandirian pangan & keberlanjutan"
        ]
    )

    # IPB Jalur Masuk
    _add_bullets_slide(
        prs,
        "IPB – Jalur Masuk",
        [
            "SNBP",
            "SNBT",
            "AFIRMASI DIKTI",
            "MANDIRI (Ketua OSIS, Talenta, SM-IPB, BUD, Kelas Internasional)",
        ]
    )

    # IPB Fakultas (split into two slides if long)
    ipb_fakultas = [
        "Fakultas Pertanian",
        "Fakultas Perikanan dan Ilmu Kelautan (FPIK)",
        "Fakultas Peternakan (FAPET)",
        "Fakultas Kehutanan dan Lingkungan (FKL)",
        "Fakultas Teknologi Pertanian (FATETA)",
        "Fakultas Matematika dan Ilmu Pengetahuan Alam (FMIPA)",
        "Fakultas Ekonomi dan Manajemen (FEM)",
        "Fakultas Ekologi Manusia (FEMA)",
        "Sekolah Kedokteran Hewan dan Biomedis (SKHB)",
        "Sekolah Bisnis (SB)",
        "Sekolah Vokasi (SV)",
    ]
    _add_bullets_slide(prs, "IPB – Fakultas & Sekolah", ipb_fakultas)

    # UI Overview
    _add_bullets_slide(
        prs,
        "Universitas Indonesia – Ringkasan",
        [
            "PTN-BH tertua dan prestisius di Indonesia",
            "Kampus utama Green Campus di Depok; Kampus Salemba di Jakarta",
            "Kampus komprehensif dan multikultural, program Vokasi hingga Doktor",
            "14 Fakultas mencakup Kesehatan, Saintek, dan Soshum",
            "Peringkat teratas nasional dengan pengakuan global",
            "Fokus: riset, inovasi, pengabdian masyarakat; lulusan berdaya saing tinggi",
        ]
    )

    # UI Jalur Masuk
    _add_bullets_slide(
        prs,
        "UI – Jalur Masuk",
        [
            "SNBP",
            "SNBT",
            "SIMAK UI",
            "Talent Scouting",
            "PPKB",
            "Seleksi Jalur Prestasi",
        ]
    )

    # UI Fakultas by rumpun
    kesehatan = [
        "Fakultas Kedokteran (FK)",
        "Fakultas Kedokteran Gigi (FKG)",
        "Fakultas Ilmu Keperawatan (FIK)",
        "Fakultas Kesehatan Masyarakat (FKM)",
        "Fakultas Farmasi (FF)",
    ]
    saintek = [
        "Fakultas Teknik (FT)",
        "Fakultas Matematika dan Ilmu Pengetahuan Alam (FMIPA)",
        "Fakultas Ilmu Komputer (Fasilkom)",
    ]
    soshum = [
        "Fakultas Hukum (FH)",
        "Fakultas Ekonomi dan Bisnis (FEB)",
        "Fakultas Ilmu Pengetahuan Budaya (FIB)",
        "Fakultas Psikologi (FPsi)",
        "Fakultas Ilmu Sosial dan Ilmu Politik (FISIP)",
        "Fakultas Ilmu Administrasi (FIA)",
    ]
    lainnya = [
        "Program Pendidikan Vokasi",
        "Sekolah Ilmu Lingkungan (SIL)",
        "Sekolah Kajian Stratejik dan Global (SKSG)",
    ]

    _add_two_column_slide(prs, "UI – Fakultas (Kesehatan & Saintek)", "Rumpun Kesehatan", kesehatan, "Rumpun Saintek", saintek)
    _add_two_column_slide(prs, "UI – Fakultas (Soshum & Program Lain)", "Rumpun Soshum", soshum, "Program/Sekolah Lain", lainnya)

    # Closing slide
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Terima kasih"
    box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(8), Inches(3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Disusun otomatis pada {datetime.now().strftime('%d %B %Y')}"
    p.font.size = Pt(18)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


@app.get("/api/ppt/ipb-ui")
def generate_ppt_ipb_ui():
    try:
        stream = build_ipb_ui_presentation()
        filename = "Profil_IPB_dan_UI.pptx"
        headers = {
            "Content-Disposition": f"attachment; filename={filename}"
        }
        return StreamingResponse(stream, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers=headers)
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
